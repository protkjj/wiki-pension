"""K-IFRS 제1019호 종업원급여 계리평가보고서 — 파워포인트(PPTX) 발표자료 생성.

report.py(엑셀 상세 보고서)와 동일한 엔진 산출값을 사용해, 고객·감사인 대상
발표용 슬라이드 덱을 만든다. 표지·핵심요약·주요가정·재무상태표 인식금액·
확정급여채무 조정·당기원가·민감도·만기구성·재직자현황·유의사항.

회사 재무자료가 필요한 항목은 disclosure_inputs로 주입하며 미제공 시 0으로 둔다.
"""

from __future__ import annotations

from datetime import date
from pathlib import Path
from typing import List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from .config import Config
from .engine import CensusResult, expected_cashflows
from .report import _BANDS, _bucketize, _class_stats, _di, _reprice

# ── 스타일 ─────────────────────────────────────────────────────────────────
FONT = "맑은 고딕"
ACCENT = RGBColor(0x1F, 0x4E, 0x79)   # 짙은 남색
INK = RGBColor(0x1F, 0x2A, 0x37)
MUTED = RGBColor(0x6B, 0x76, 0x84)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
HEAD_FILL = RGBColor(0xDC, 0xE6, 0xF1)
TOT_FILL = RGBColor(0xEE, 0xF3, 0xF8)
BAND = RGBColor(0xF5, 0xF7, 0xFA)

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _won(n) -> str:
    try:
        return f"{round(n):,}"
    except Exception:
        return str(n)


def _pct(x) -> str:
    return f"{x * 100:.3f}%"


def _txt(tf, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = FONT
    r.font.color.rgb = color
    return p


def _add_textbox(slide, left, top, width, height, text, **kw):
    box = slide.shapes.add_textbox(left, top, width, height)
    _txt(box.text_frame, text, **kw)
    return box


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _header(slide, title, subtitle=None):
    """상단 남색 제목 바."""
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _txt(tf, title, size=22, bold=True, color=WHITE)
    if subtitle:
        _add_textbox(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.4),
                     subtitle, size=11, color=MUTED)


def _add_table(slide, headers, rows, top, left=Inches(0.6), width=None,
               col_aligns=None, total_last=False, font_size=13):
    """헤더 1행 + 데이터 n행 표. rows: list[list[str]]."""
    width = width or (SLIDE_W - Inches(1.2))
    ncol = len(headers)
    nrow = len(rows) + 1
    height = Inches(0.42) * nrow
    gtbl = slide.shapes.add_table(nrow, ncol, left, top, width, height)
    tbl = gtbl.table
    col_aligns = col_aligns or ([PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * (ncol - 1))

    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = HEAD_FILL
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
        tfr = cell.text_frame
        _txt(tfr, str(h), size=font_size, bold=True, color=INK, align=PP_ALIGN.CENTER)

    for i, row in enumerate(rows, start=1):
        is_tot = total_last and i == len(rows)
        for j, v in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TOT_FILL if is_tot else (BAND if i % 2 == 0 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_top = Pt(1); cell.margin_bottom = Pt(1)
            _txt(cell.text_frame, str(v), size=font_size, bold=is_tot,
                 color=INK, align=col_aligns[j])
    return gtbl


def _metric_tile(slide, left, top, w, h, label, value, sub=None):
    card = slide.shapes.add_shape(1, left, top, w, h)
    card.fill.solid(); card.fill.fore_color.rgb = TOT_FILL
    card.line.color.rgb = HEAD_FILL
    tf = card.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    _txt(tf, label, size=13, color=MUTED)
    p = tf.add_paragraph()
    r = p.add_run(); r.text = value
    r.font.size = Pt(26); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = ACCENT
    if sub:
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(11); r2.font.name = FONT; r2.font.color.rgb = MUTED


# ═══════════════════════════════════════════════════════════════════════════
def write_report_pptx(
    out_path,
    company: str,
    valuation_date: date,
    config: Config,
    census_result: CensusResult,
    records,
    tables,
    disclosure_inputs: Optional[dict] = None,
    plan_info: Optional[dict] = None,
    prior: Optional[dict] = None,
    valuer: str = "보험계리사",
) -> Path:
    di = disclosure_inputs or {}
    y = valuation_date.year
    dbo = round(census_result.total_dbo)
    csc = round(census_result.total_csc)
    disc = config.discount_rate.flat
    sal = config.salary_increase_rate.flat
    ret_age = config.retirement_age.default

    cf = expected_cashflows(records, config, tables)
    dur = (cf["연도"] * cf["현재가치"]).sum() / cf["현재가치"].sum() if len(cf) and cf["현재가치"].sum() else 0.0
    n_calc = len(census_result.results)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── 1) 표지 ──
    s = _blank(prs)
    band = s.shapes.add_shape(1, 0, Inches(2.4), SLIDE_W, Inches(2.7))
    band.fill.solid(); band.fill.fore_color.rgb = ACCENT; band.line.fill.background()
    _add_textbox(s, Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.6),
                 f"{company} 를 위한", size=18, color=WHITE)
    _add_textbox(s, Inches(0.8), Inches(3.25), Inches(11.7), Inches(1.0),
                 "K-IFRS 제1019호 종업원급여 계리평가보고서", size=30, bold=True, color=WHITE)
    _add_textbox(s, Inches(0.8), Inches(5.4), Inches(11.7), Inches(0.5),
                 f"산출기준일: {y}년 {valuation_date.month}월 {valuation_date.day}일",
                 size=16, bold=True, color=INK)
    _add_textbox(s, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.4),
                 f"작성: {valuer}", size=13, color=MUTED)

    # ── 2) 핵심 요약 ──
    s = _blank(prs)
    _header(s, "핵심 요약", "K-IFRS 제1019호 계리평가 결과 요약")
    tiles = [
        ("확정급여채무(DBO)", f"{_won(dbo)} 원", "재무상태표상 현재가치"),
        ("당기근무원가(CSC)", f"{_won(csc)} 원", "당기손익 인식분"),
        ("계산대상 인원", f"{n_calc:,} 명", f"제외 {len(census_result.excluded_emp_ids):,} 명"),
        ("가중평균만기", f"{dur:.2f} 년", "확정급여채무 듀레이션"),
    ]
    x0, gap, w, h = Inches(0.6), Inches(0.25), Inches(2.95), Inches(1.7)
    for i, (lb, val, sub) in enumerate(tiles):
        _metric_tile(s, x0 + i * (w + gap), Inches(1.7), w, h, lb, val, sub)
    _add_textbox(s, Inches(0.6), Inches(3.9), Inches(12), Inches(2.6),
                 "· 예측단위적립방식(PUC)으로 정년까지의 예상 퇴직급여를 현재 근속에 비례해 배분·할인하여 산정.\n"
                 "· 임원·계약직·정년초과자는 회사 제공 금액(간편법)을 준용.\n"
                 "· 상세 수치는 첨부 엑셀 계리평가보고서(주석공시사항·개인별명세)를 참조.",
                 size=14, color=INK)

    # ── 3) 주요 보험수리적 가정 ──
    s = _blank(prs)
    _header(s, "주요 보험수리적 가정")
    rows = [
        ["할인율", _pct(disc)],
        ["임금상승률(Base-up)", _pct(sal)],
        ["퇴직률", f"{config.retirement_rate_basis} 기준 참조율"],
        ["사망률", "성별·연령별 참조율"],
        ["정년", f"만 {ret_age} 세"],
        ["기대잔존근무년수(가중평균만기)", f"{dur:.2f} 년"],
    ]
    _add_table(s, ["가정 항목", "적용값"], rows, Inches(1.6),
               width=Inches(9.0), col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT])

    # ── 4) 재무상태표에 인식된 순확정급여부채(자산) ──
    assets = _di(di, "plan_assets"); npc = _di(di, "npc_conversion"); ceiling = _di(di, "asset_ceiling")
    net = round(-assets + dbo + ceiling + npc)
    s = _blank(prs)
    _header(s, "재무상태표에 인식된 순확정급여부채(자산)", "(단위: 원)")
    rows = [
        ["확정급여채무의 현재가치", _won(dbo)],
        ["사외적립자산", f"({_won(assets)})"],
        ["자산인식상한효과", _won(ceiling)],
        ["국민연금전환금", f"({_won(npc)})"],
        ["순확정급여부채(자산)", _won(net)],
    ]
    _add_table(s, ["항목", "당기말"], rows, Inches(1.6),
               width=Inches(9.0), col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT], total_last=True)

    # ── 5) 확정급여채무의 현재가치 조정내역 ──
    ic = round(_di(di, "dbo_begin") * disc) if _di(di, "dbo_begin") else _di(di, "interest_cost")
    s = _blank(prs)
    _header(s, "확정급여채무의 현재가치 조정내역", "(단위: 원)")
    rows = [
        ["기초 확정급여채무", _won(_di(di, "dbo_begin"))],
        ["당기근무원가", _won(csc)],
        ["이자비용", _won(ic)],
        ["급여지급액(중간정산 포함)", f"({_won(_di(di, 'benefits_paid_dbo'))})"],
        ["재측정손익(가정·경험 조정)", _won(_di(di, "remeasure_demographic")
                                    + _di(di, "remeasure_financial") + _di(di, "remeasure_experience"))],
        ["기말 확정급여채무", _won(dbo)],
    ]
    _add_table(s, ["항목", "당기"], rows, Inches(1.6),
               width=Inches(9.5), col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT], total_last=True)

    # ── 6) 당기손익으로 인식하는 확정급여원가 ──
    net_interest = _di(di, "net_interest")
    s = _blank(prs)
    _header(s, "당기손익으로 인식하는 확정급여원가", "(단위: 원)")
    rows = [
        ["당기근무원가", _won(csc)],
        ["과거근무원가와 정산 손익", _won(0)],
        ["순확정급여부채(자산)의 순이자", _won(net_interest)],
        ["당기손익 인식 확정급여원가", _won(round(csc + net_interest))],
    ]
    _add_table(s, ["항목", "당기"], rows, Inches(1.6),
               width=Inches(9.0), col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT], total_last=True)

    # ── 7) 민감도 분석 ──
    d_up, d_dn = _reprice(records, config, tables, dd=+0.01), _reprice(records, config, tables, dd=-0.01)
    s_up, s_dn = _reprice(records, config, tables, sd=+0.01), _reprice(records, config, tables, sd=-0.01)
    s = _blank(prs)
    _header(s, "민감도 분석 — 확정급여채무 변동", "주요 가정 ±1.00%p (단위: 원)")
    rows = [
        ["할인율", _won(round(d_up)), _won(dbo), _won(round(d_dn))],
        ["임금상승률(Base-up)", _won(round(s_up)), _won(dbo), _won(round(s_dn))],
    ]
    _add_table(s, ["가정", "1.00%p 상승", "기준", "1.00%p 하락"], rows, Inches(1.7),
               width=Inches(11.5),
               col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT, PP_ALIGN.RIGHT])
    _add_textbox(s, Inches(0.6), Inches(3.4), Inches(12), Inches(1.0),
                 "· 할인율이 낮아지면(하락) 확정급여채무는 증가하고, 임금상승률이 높아지면 채무가 증가합니다.",
                 size=13, color=MUTED)

    # ── 8) 만기구성정보 ──
    pv_b = _bucketize(cf, "현재가치")
    simple_dbo = census_result.subtotal_by_plan.get(2, {}).get("DBO", 0.0)
    s = _blank(prs)
    _header(s, "확정급여채무의 만기구성정보", f"가중평균만기 {dur:.2f}년 (단위: 원)")
    rows = [[name, _won(round(pv_b[name]))] for name in _BANDS]
    if simple_dbo:
        rows.append(["간편법 평가분(제도2)", _won(round(simple_dbo))])
    rows.append(["합계", _won(dbo)])
    _add_table(s, ["만기", "확정급여채무"], rows, Inches(1.55),
               left=Inches(2.4), width=Inches(8.5),
               col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.RIGHT], total_last=True, font_size=11)

    # ── 9) 재직자 현황 ──
    stats, total = _class_stats(census_result)
    s = _blank(prs)
    _header(s, f"{valuation_date} 현재 재직자 현황")
    rows = [[name, f"{n:,}", f"{aage:.1f}", f"{asvc:.1f}", _won(round(ssal))]
            for name, n, aage, asvc, ssal in stats]
    rows.append(["합계", f"{total[1]:,}", f"{total[2]:.1f}", f"{total[3]:.1f}", _won(round(total[4]))])
    _add_table(s, ["구분", "인원(명)", "평균연령", "평균근속(년)", "평균임금 합(원)"], rows, Inches(1.7),
               width=Inches(11.0),
               col_aligns=[PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * 4, total_last=True)

    # ── 10) 유의사항 ──
    s = _blank(prs)
    _header(s, "유의사항")
    _add_textbox(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.0),
                 "· 본 자료는 K-IFRS 제1019호 및 한국계리업무기준 제3편에 따라 작성되었습니다.\n\n"
                 "· 회계 결산 지원 및 외부감사 목적으로 활용하며, 세무조정·부담금 산출 등 다른 목적에는 "
                 "부적절할 수 있습니다.\n\n"
                 "· 사외적립자산·기여금·전기 대비 조정 및 재측정손익은 회사가 제공한 재무자료에 기초하며, "
                 "미제공 항목은 0으로 표시됩니다.\n\n"
                 "· 상세 수치(주석공시사항 8개 표·개인별명세)는 첨부 엑셀 계리평가보고서를 참조하십시오.",
                 size=15, color=INK)

    prs.save(str(Path(out_path)))
    return Path(out_path)
